# ──────────────────────────────────────────────
# INDEX — Script de indexação de documentos no Pinecone
# ──────────────────────────────────────────────

"""
Script para indexar documentos no Pinecone.
Execute: python index.py

Este script:
1. Carrega documentos da pasta 'documentos/'
2. Extrai texto de múltiplos formatos (PDF, DOCX, XLSX, PPTX, CSV, HTML, MD, TXT)
3. Limpa e divide em chunks
4. Gera embeddings via API (com fallback local)
4. Envia para o Pinecone
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path
from typing import List

from tqdm import tqdm

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import get_config
from embeddings.provider import get_embedding_provider
from vectorstore.pinecone_client import pinecone_client

# ──────────────────────────────────────────────
# CONFIGURAÇÃO
# ──────────────────────────────────────────────

BASE_DIR = Path(__file__).parent
config = get_config()

PASTA_DOCUMENTOS = BASE_DIR / "documentos"
CHUNK_SIZE = config.chunking.chunk_size
CHUNK_OVERLAP = config.chunking.chunk_overlap
MODELO_EMBEDDING = config.embedding.model
DIMENSION = config.embedding.dimension

# ──────────────────────────────────────────────
# CURADORIA — Filtros de qualidade
# ──────────────────────────────────────────────

CURADORIA = config.curadoria
IGNORAR_SE_CONTER = [p.lower() for p in CURADORIA.ignorar_se_conter]
IGNORAR_EXATOS = [p.lower() for p in CURADORIA.ignorar_exatos]
MANTER_VERSAO_OFICIAL = CURADORIA.manter_versao_oficial

_PADRAO_VERSAO = re.compile(r"[-_ ]v?\d+(\.\d+)*$", re.IGNORECASE)


def _arquivo_ignorado(caminho: Path) -> bool:
    """Retorna True se o arquivo deve ser ignorado (rascunhos, backups, etc.)."""
    nome = caminho.stem.lower()
    nome_completo = caminho.name.lower()

    if nome_completo in IGNORAR_EXATOS:
        return True

    for padrao in IGNORAR_SE_CONTER:
        if padrao in nome:
            return True

    return False


def _filtrar_versoes(arquivos: List[Path]) -> List[Path]:
    """Mantém apenas a versão oficial (sem sufixo) quando há múltiplas do mesmo doc."""
    if not MANTER_VERSAO_OFICIAL:
        return arquivos

    grupos = {}
    for arq in arquivos:
        base = _PADRAO_VERSAO.sub("", arq.stem).lower()
        grupos.setdefault(base, []).append(arq)

    mantidos = set()
    for base, grupo in grupos.items():
        if len(grupo) == 1:
            mantidos.add(grupo[0])
            continue

        sem_versao = [a for a in grupo if not _PADRAO_VERSAO.search(a.stem)]
        if sem_versao:
            mantidos.add(sem_versao[0])
            ignorados = [a for a in grupo if a != sem_versao[0]]
        else:
            grupo_ordenado = sorted(grupo, key=lambda a: a.stat().st_mtime, reverse=True)
            mantidos.add(grupo_ordenado[0])
            ignorados = grupo_ordenado[1:]

        for ign in ignorados:
            print(f"  Curadoria: ignorado {ign.name} (versao anterior de {base})")

    return list(mantidos)


# ──────────────────────────────────────────────
# EXTRATORES DE TEXTO POR FORMATO
# ──────────────────────────────────────────────

def extrair_texto_pdf(caminho: Path) -> str:
    """Extrai texto de PDF via PyMuPDF."""
    import fitz
    doc = fitz.open(str(caminho))
    try:
        return "\n".join(pagina.get_text() for pagina in doc)
    finally:
        doc.close()


def extrair_texto_docx(caminho: Path) -> str:
    """Extrai texto de arquivos .docx."""
    from docx import Document as DocxDocument
    doc = DocxDocument(str(caminho))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def extrair_texto_xlsx(caminho: Path) -> str:
    """Extrai texto de planilhas .xlsx, células separadas por pipe."""
    import openpyxl
    wb = openpyxl.load_workbook(str(caminho), data_only=True)
    try:
        linhas = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                celulas = [str(c) for c in row if c is not None]
                if celulas:
                    linhas.append(" | ".join(celulas))
        return "\n".join(linhas)
    finally:
        wb.close()


def extrair_texto_pptx(caminho: Path) -> str:
    """Extrai texto de apresentações .pptx."""
    from pptx import Presentation
    prs = Presentation(str(caminho))
    try:
        texto = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texto.append(shape.text)
        return "\n".join(texto)
    finally:
        if hasattr(prs, "close"):
            prs.close()


def _extrair_tabelas_html(soup) -> str:
    """Converte tabelas HTML para formato pipe (|) legível, como markdown."""
    linhas = []
    for table in soup.find_all("table"):
        cabecalho = []
        for th in table.find_all("th"):
            cabecalho.append(th.get_text(strip=True))
        if cabecalho:
            linhas.append("| " + " | ".join(cabecalho) + " |")
            linhas.append("|-" + "-|-".join("---" for _ in cabecalho) + "-|")
        for tr in table.find_all("tr"):
            celulas = []
            for td in tr.find_all("td"):
                celulas.append(td.get_text(strip=True))
            if celulas:
                linhas.append("| " + " | ".join(celulas) + " |")
        linhas.append("")
    return "\n".join(linhas)


def extrair_texto_html(caminho: Path) -> str:
    """Extrai texto de HTML, convertendo tabelas para pipe e removendo-as do texto plano."""
    from bs4 import BeautifulSoup
    with open(caminho, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    tabelas = _extrair_tabelas_html(soup)
    for table in soup.find_all("table"):
        table.decompose()
    texto = soup.get_text(separator="\n")
    if tabelas:
        return tabelas + "\n\n" + texto
    return texto


def extrair_texto_csv(caminho: Path) -> str:
    """Extrai texto de CSV com detecção automática de delimitador e fallback de encoding."""
    with open(caminho, "r", encoding="utf-8") as f:
        sample = f.read(2048)
    try:
        delimiter = csv.Sniffer().sniff(sample).delimiter
    except csv.Error:
        delimiter = ","
    try:
        import pandas as pd
        df = pd.read_csv(str(caminho), encoding="utf-8", sep=delimiter)
    except UnicodeDecodeError:
        df = pd.read_csv(str(caminho), encoding="latin-1", sep=delimiter)
    return df.to_string(index=False, na_rep="")


def extrair_texto_json(caminho: Path) -> str:
    """Extrai texto de JSON com validação de formato."""
    with open(caminho, "r", encoding="utf-8") as f:
        try:
            dados = json.load(f)
        except json.JSONDecodeError as e:
            raise ValueError(f"JSON inválido: {e}")
    return json.dumps(dados, ensure_ascii=False, indent=2)


def extrair_texto_md(caminho: Path) -> str:
    """Extrai texto de markdown/txt com fallback de encoding."""
    try:
        with open(caminho, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(caminho, "r", encoding="latin-1") as f:
            return f.read()


# Mapeamento extensão -> extrator
EXTRATORES = {
    ".pdf": extrair_texto_pdf,
    ".docx": extrair_texto_docx,
    ".xlsx": extrair_texto_xlsx,
    ".pptx": extrair_texto_pptx,
    ".html": extrair_texto_html,
    ".htm": extrair_texto_html,
    ".csv": extrair_texto_csv,
    ".json": extrair_texto_json,
    ".md": extrair_texto_md,
    ".txt": extrair_texto_md,
}

# ──────────────────────────────────────────────
# LIMPEZA DE RUÍDOS — remove cabeçalhos, rodapés, numeração de página
# ──────────────────────────────────────────────

def _limpar_texto(texto: str) -> str:
    """Remove ruídos comuns (cabeçalhos, rodapés, paginação, linhas repetitivas)."""
    if not texto or not texto.strip():
        return ""
    linhas = texto.split("\n")
    limpas = []
    contagem = {}
    for linha in linhas:
        chave = linha.strip().lower()
        if len(chave) > 3:
            contagem[chave] = contagem.get(chave, 0) + 1
    total = len(linhas)
    limiar = max(4, total // 20)
    for linha in linhas:
        s = linha.strip()
        if not s:
            limpas.append("")
            continue
        if re.match(r"^\d+\s*/\s*\d+$", s):
            continue
        if re.match(r"^p[aá]gina\s+\d+", s, re.IGNORECASE):
            continue
        if re.match(r"^p[aá]g\.?\s*\d+", s, re.IGNORECASE):
            continue
        if re.match(r"^-\s*\d+\s*-$", s):
            continue
        if re.match(r"^confidencial$", s, re.IGNORECASE):
            continue
        if re.match(r"^documento\s+interno", s, re.IGNORECASE):
            continue
        chave = s.lower()
        if len(chave) > 3 and contagem.get(chave, 0) >= limiar:
            continue
        limpas.append(s)
    resultado = "\n".join(limpas)
    resultado = re.sub(r"\n{3,}", "\n\n", resultado)
    return resultado.strip()


# ──────────────────────────────────────────────
# FUNÇÃO PRINCIPAL — processa e indexa documentos
# ──────────────────────────────────────────────

def processar_documentos() -> bool:
    """
    Percorre as pastas de departamento dentro de documentos/,
    extrai o texto de cada arquivo suportado, divide em chunks
    e indexa tudo no Pinecone para consulta via busca semântica.
    """

    if not PASTA_DOCUMENTOS.is_dir():
        print(f"Pasta {PASTA_DOCUMENTOS} não encontrada ou não é um diretório. Crie a estrutura:")
        print("  documentos/")
        print("    rh/")
        print("    financeiro/")
        print("    juridico/")
        return False

    if pinecone_client._config:
        # Força recriação do índice se necessário
        pass

    # Remove índice anterior se existir (para reindexar limpo)
    # Nota: Em produção, pode querer fazer upsert incremental
    print("Removendo índice anterior (se existir)...")
    try:
        pc = pinecone_client.get_client()
        indexes = pc.list_indexes().names()
        if config.pinecone.index_name in indexes:
            pc.delete_index(config.pinecone.index_name)
            print("Índice anterior removido.")
    except Exception as e:
        print(f"Aviso ao remover índice anterior: {e}")

    # Carrega embeddings
    print(f"Carregando modelo de embedding: {MODELO_EMBEDDING}")
    try:
        embeddings = get_embedding_provider(model=MODELO_EMBEDDING)
    except Exception as e:
        print(f"Erro ao carregar modelo de embedding '{MODELO_EMBEDDING}': {e}")
        return False

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ".", " ", ""],
    )

    todos_documentos = []

    try:
        pastas_departamento = [p for p in PASTA_DOCUMENTOS.iterdir() if p.is_dir()]
    except PermissionError:
        print(f"Erro: sem permissão para ler a pasta {PASTA_DOCUMENTOS}.")
        return False

    if not pastas_departamento:
        print("Nenhuma subpasta encontrada em documentos/.")
        print("Crie pastas como: documentos/rh, documentos/financeiro, documentos/juridico")
        return False

    for pasta_dep in pastas_departamento:
        departamento = pasta_dep.name
        print(f"\nProcessando departamento: {departamento}")

        arquivos = set()
        for ext in EXTRATORES:
            arquivos.update(pasta_dep.glob(f"**/*{ext}"))
            arquivos.update(pasta_dep.glob(f"**/*{ext.upper()}"))
            arquivos.update(pasta_dep.glob(f"**/*{ext.capitalize()}"))

        if not arquivos:
            print(f"  Nenhum arquivo suportado encontrado em {pasta_dep}")
            continue

        arquivos = _filtrar_versoes(arquivos)
        arquivos_ordenados = sorted(arquivos, key=lambda a: a.name)
        if arquivos_ordenados:
            print(f"  Processando {len(arquivos_ordenados)} arquivos...")

        for caminho in tqdm(arquivos_ordenados, desc=f"  {departamento}", unit="arquivo", leave=False):
            if _arquivo_ignorado(caminho):
                print(f"  Curadoria: ignorado {caminho.name}")
                continue
            extrator = EXTRATORES.get(caminho.suffix.lower())
            if not extrator:
                continue

            try:
                texto = extrator(str(caminho))
                texto = _limpar_texto(texto)
                if not texto.strip():
                    print(f"  WARN Vazio: {caminho.name}")
                    continue

                chunks = splitter.split_text(texto)

                metadados_base = {
                    "departamento": departamento,
                    "arquivo": str(caminho.relative_to(PASTA_DOCUMENTOS)),
                    "fonte": caminho.name,
                }

                for i, chunk in enumerate(chunks):
                    metadados = metadados_base.copy()
                    metadados["chunk_id"] = i
                    todos_documentos.append(
                        Document(page_content=chunk, metadata=metadados)
                    )

                print(f"  OK {caminho.name} -> {len(chunks)} chunks")

            except Exception as e:
                print(f"  ERR {caminho.name}: {e}")

    if not todos_documentos:
        print("\nNenhum chunk gerado. Verifique os documentos.")
        return False

    # Indexa no Pinecone
    print(f"\nGerando embeddings e indexando {len(todos_documentos)} chunks no Pinecone...")
    try:
        pinecone_client.add_documents(
            documents=todos_documentos,
            embeddings=embeddings,
            namespace="default",
        )
    except Exception as e:
        print(f"Erro ao indexar documentos no Pinecone: {e}")
        return False

    # Conta arquivos únicos
    arquivos_processados = len(set(
        doc.metadata.get("arquivo", doc.metadata.get("fonte", ""))
        for doc in todos_documentos
    ))

    print(f"\n✅ Concluído! {len(todos_documentos)} chunks de {arquivos_processados} arquivos indexados no Pinecone")
    print(f"  Índice: {config.pinecone.index_name}")
    print(f"  Departamentos: {[p.name for p in pastas_departamento]}")
    return True


if __name__ == "__main__":
    processar_documentos()