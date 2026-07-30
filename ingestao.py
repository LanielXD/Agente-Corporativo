# ──────────────────────────────────────────────
# IMPORTAÇÕES
# ──────────────────────────────────────────────

import shutil
import csv
import json
import re
from pathlib import Path

import yaml
import logger
from tqdm import tqdm

BASE_DIR = Path(__file__).parent
import fitz
import openpyxl
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document

# ──────────────────────────────────────────────
# CONFIGURAÇÃO
# ──────────────────────────────────────────────

try:
    with open(BASE_DIR / "config.yaml", "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)
except FileNotFoundError as e:
    raise RuntimeError("Arquivo config.yaml não encontrado.") from e
except yaml.YAMLError as e:
    raise RuntimeError(f"Erro ao ler config.yaml: {e}") from e

if config is None:
    raise RuntimeError("config.yaml está vazio.")

PASTA_DOCUMENTOS = BASE_DIR / "documentos"
PASTA_CHROMA = BASE_DIR / "chroma_db"
CHUNK_TAMANHO = config.get("chunk_tamanho", 500)
CHUNK_SOBREPOSICAO = config.get("chunk_sobreposicao", 50)
MODELO_EMBEDDING = config.get("modelo_embedding")
if not MODELO_EMBEDDING:
    raise RuntimeError("'modelo_embedding' não definido em config.yaml.")

# ──────────────────────────────────────────────
# CURADORIA — filtros de qualidade
# ──────────────────────────────────────────────

CURADORIA = config.get("curadoria", {})
IGNORAR_SE_CONTER = [p.lower() for p in CURADORIA.get("ignorar_se_conter", [])]
IGNORAR_EXATOS = [p.lower() for p in CURADORIA.get("ignorar_exatos", [])]
MANTER_VERSAO_OFICIAL = CURADORIA.get("manter_versao_oficial", False)

# Expressão para detectar sufixos de versão no nome do arquivo
_PADRAO_VERSAO = re.compile(r"[-_ ]v?\d+(\.\d+)*$", re.IGNORECASE)


def _arquivo_ignorado(caminho):
    """Retorna True se o arquivo deve ser ignorado (rascunhos, backups, etc.)."""
    nome = caminho.stem.lower()
    nome_completo = caminho.name.lower()

    if nome_completo in IGNORAR_EXATOS:
        return True

    for padrao in IGNORAR_SE_CONTER:
        if padrao in nome:
            return True

    return False


def _filtrar_versoes(arquivos, log_fn=print):
    """Mantém apenas a versão oficial (sem sufixo) quando há múltiplas do mesmo doc."""
    if not MANTER_VERSAO_OFICIAL:
        return arquivos

    # Agrupa por nome base (remove sufixo de versão)
    grupos = {}
    for arq in arquivos:
        base = _PADRAO_VERSAO.sub("", arq.stem).lower()
        grupos.setdefault(base, []).append(arq)

    # Para cada grupo com mais de um arquivo, mantém apenas o sem versão ou o mais recente
    mantidos = set()
    for base, grupo in grupos.items():
        if len(grupo) == 1:
            mantidos.add(grupo[0])
            continue

        sem_versao = [a for a in grupo if not _PADRAO_VERSAO.search(a.stem)]
        if sem_versao:
            mantidos.add(sem_versao[0])
            ignorados = [a for a in grupo if a != sem_versao[0]]
        else:
            grupo_ordenado = sorted(grupo, key=lambda a: a.stat().st_mtime, reverse=True)
            mantidos.add(grupo_ordenado[0])
            ignorados = grupo_ordenado[1:]

        for ign in ignorados:
            log_fn(f"  Curadoria: ignorado {ign.name} (versao anterior de {base})")

    return list(mantidos)


# ──────────────────────────────────────────────
# FUNÇÕES DE EXTRAÇÃO DE TEXTO POR FORMATO
# ──────────────────────────────────────────────

def extrair_texto_pdf(caminho):
    """Extrai texto de PDF via PyMuPDF."""
    doc = fitz.open(caminho)
    try:
        return "\n".join(pagina.get_text() for pagina in doc)
    finally:
        doc.close()


def extrair_texto_docx(caminho):
    """Extrai texto de arquivos .docx."""
    doc = DocxDocument(caminho)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def extrair_texto_xlsx(caminho):
    """Extrai texto de planilhas .xlsx, células separadas por pipe."""
    wb = openpyxl.load_workbook(caminho, data_only=True)
    try:
        linhas = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                celulas = [str(c) for c in row if c is not None]
                if celulas:
                    linhas.append(" | ".join(celulas))
        return "\n".join(linhas)
    finally:
        wb.close()


def extrair_texto_pptx(caminho):
    """Extrai texto de apresentações .pptx."""
    prs = Presentation(caminho)
    try:
        texto = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texto.append(shape.text)
        return "\n".join(texto)
    finally:
        if hasattr(prs, "close"):
            prs.close()


def _extrair_tabelas_html(soup):
    """Converte tabelas HTML para formato pipe (|) legível, como markdown."""
    linhas = []
    for table in soup.find_all("table"):
        cabecalho = []
        for th in table.find_all("th"):
            cabecalho.append(th.get_text(strip=True))
        if cabecalho:
            linhas.append("| " + " | ".join(cabecalho) + " |")
            linhas.append("|-" + "-|-".join("---" for _ in cabecalho) + "-|")
        for tr in table.find_all("tr"):
            celulas = []
            for td in tr.find_all("td"):
                celulas.append(td.get_text(strip=True))
            if celulas:
                linhas.append("| " + " | ".join(celulas) + " |")
        linhas.append("")
    return "\n".join(linhas)


def extrair_texto_html(caminho):
    """Extrai texto de HTML, convertendo tabelas para pipe e removendo-as do texto plano."""
    with open(caminho, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    tabelas = _extrair_tabelas_html(soup)
    for table in soup.find_all("table"):
        table.decompose()
    texto = soup.get_text(separator="\n")
    if tabelas:
        return tabelas + "\n\n" + texto
    return texto


def extrair_texto_csv(caminho):
    """Extrai texto de CSV com detecção automática de delimitador e fallback de encoding."""
    with open(caminho, "r", encoding="utf-8") as f:
        sample = f.read(2048)
    try:
        delimiter = csv.Sniffer().sniff(sample).delimiter
    except csv.Error:
        delimiter = ","
    try:
        df = pd.read_csv(caminho, encoding="utf-8", sep=delimiter)
    except UnicodeDecodeError:
        df = pd.read_csv(caminho, encoding="latin-1", sep=delimiter)
    return df.to_string(index=False, na_rep="")


def extrair_texto_json(caminho):
    """Extrai texto de JSON com validação de formato."""
    with open(caminho, "r", encoding="utf-8") as f:
        try:
            dados = json.load(f)
        except json.JSONDecodeError as e:
            raise ValueError(f"JSON inválido: {e}")
    return json.dumps(dados, ensure_ascii=False, indent=2)


def extrair_texto_md(caminho):
    """Extrai texto de markdown/txt com fallback de encoding."""
    try:
        with open(caminho, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(caminho, "r", encoding="latin-1") as f:
            return f.read()


# Mapeamento entre extensão de arquivo e função extratora
EXTRATORES = {
    ".pdf": extrair_texto_pdf,
    ".docx": extrair_texto_docx,
    ".xlsx": extrair_texto_xlsx,
    ".pptx": extrair_texto_pptx,
    ".html": extrair_texto_html,
    ".htm": extrair_texto_html,
    ".csv": extrair_texto_csv,
    ".json": extrair_texto_json,
    ".md": extrair_texto_md,
    ".txt": extrair_texto_md,
}

# ──────────────────────────────────────────────
# LIMPEZA DE RUÍDOS — remove cabeçalhos, rodapés, numeração de página
# ──────────────────────────────────────────────


def _limpar_texto(texto):
    """Remove ruídos comuns (cabeçalhos, rodapés, paginação, linhas repetitivas)."""
    if not texto or not texto.strip():
        return ""
    linhas = texto.split("\n")
    limpas = []
    contagem = {}
    for linha in linhas:
        chave = linha.strip().lower()
        if len(chave) > 3:
            contagem[chave] = contagem.get(chave, 0) + 1
    total = len(linhas)
    limiar = max(4, total // 20)
    for linha in linhas:
        s = linha.strip()
        if not s:
            limpas.append("")
            continue
        if re.match(r"^\d+\s*/\s*\d+$", s):
            continue
        if re.match(r"^p[aá]gina\s+\d+", s, re.IGNORECASE):
            continue
        if re.match(r"^p[aá]g\.?\s*\d+", s, re.IGNORECASE):
            continue
        if re.match(r"^-\s*\d+\s*-$", s):
            continue
        if re.match(r"^confidencial$", s, re.IGNORECASE):
            continue
        if re.match(r"^documento\s+interno", s, re.IGNORECASE):
            continue
        chave = s.lower()
        if len(chave) > 3 and contagem.get(chave, 0) >= limiar:
            continue
        limpas.append(s)
    resultado = "\n".join(limpas)
    resultado = re.sub(r"\n{3,}", "\n\n", resultado)
    return resultado.strip()


# ──────────────────────────────────────────────
# FUNÇÃO PRINCIPAL — processa e indexa documentos
# ──────────────────────────────────────────────

def processar_documentos(log_fn=print):
    """
    Percorre as pastas de departamento dentro de documentos/,
    extrai o texto de cada arquivo suportado, divide em chunks
    e indexa tudo no ChromaDB para consulta via busca semântica.
    
    Args:
        log_fn: Função de log (padrão: print). Use lambda *a, **k: None para silenciar.
    """
    if not PASTA_DOCUMENTOS.is_dir():
        log_fn(f"Pasta {PASTA_DOCUMENTOS} não encontrada ou não é um diretório. Crie a estrutura:")
        log_fn("  documentos/")
        log_fn("    rh/")
        log_fn("    financeiro/")
        log_fn("    juridico/")
        return False

    if PASTA_CHROMA.exists():
        shutil.rmtree(PASTA_CHROMA)
        log_fn("Banco anterior removido. Reindexando...")

    try:
        embeddings = HuggingFaceEmbeddings(model_name=MODELO_EMBEDDING)
    except Exception as e:
        log_fn(f"Erro ao carregar modelo de embedding '{MODELO_EMBEDDING}': {e}")
        return False

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_TAMANHO,
        chunk_overlap=CHUNK_SOBREPOSICAO,
        separators=["\n\n", "\n", ".", " ", ""],
    )

    todos_documentos = []

    try:
        pastas_departamento = [p for p in PASTA_DOCUMENTOS.iterdir() if p.is_dir()]
    except PermissionError:
        log_fn(f"Erro: sem permissão para ler a pasta {PASTA_DOCUMENTOS}.")
        return False
    if not pastas_departamento:
        log_fn("Nenhuma subpasta encontrada em documentos/.")
        log_fn("Crie pastas como: documentos/rh, documentos/financeiro, documentos/juridico")
        return False

    for pasta_dep in pastas_departamento:
        departamento = pasta_dep.name
        log_fn(f"\nProcessando departamento: {departamento}")

        arquivos = set()
        for ext in EXTRATORES:
            arquivos.update(pasta_dep.glob(f"**/*{ext}"))
            arquivos.update(pasta_dep.glob(f"**/*{ext.upper()}"))
            arquivos.update(pasta_dep.glob(f"**/*{ext.capitalize()}"))

        if not arquivos:
            log_fn(f"  Nenhum arquivo suportado encontrado em {pasta_dep}")
            continue

        arquivos = _filtrar_versoes(arquivos, log_fn)
        arquivos_ordenados = sorted(arquivos, key=lambda a: a.name)
        if arquivos_ordenados:
            log_fn(f"  Processando {len(arquivos_ordenados)} arquivos...")
        for caminho in tqdm(arquivos_ordenados, desc=f"  {departamento}", unit="arquivo", leave=False):
            if _arquivo_ignorado(caminho):
                log_fn(f"  Curadoria: ignorado {caminho.name}")
                continue
            extrator = EXTRATORES.get(caminho.suffix.lower())
            if not extrator:
                continue

            try:
                texto = extrator(str(caminho))
                texto = _limpar_texto(texto)
                if not texto.strip():
                    log_fn(f"  WARN Vazio: {caminho.name}")
                    continue

                chunks = splitter.split_text(texto)

                metadados_base = {
                    "departamento": departamento,
                    "arquivo": str(caminho.relative_to(PASTA_DOCUMENTOS)),
                    "fonte": caminho.name,
                }

                for i, chunk in enumerate(chunks):
                    metadados = metadados_base.copy()
                    metadados["chunk_id"] = i
                    todos_documentos.append(
                        Document(page_content=chunk, metadata=metadados)
                    )

                log_fn(f"  OK {caminho.name} -> {len(chunks)} chunks")

            except Exception as e:
                log_fn(f"  ERR {caminho.name}: {e}")

    if not todos_documentos:
        log_fn("\nNenhum chunk gerado. Verifique os documentos.")
        logger.ingestao(0, 0, False)
        return False

    # Conta arquivos que geraram pelo menos 1 chunk
    arquivos_processados = len(set(
        doc.metadata.get("arquivo", doc.metadata.get("fonte", ""))
        for doc in todos_documentos
    ))

    log_fn(f"\nGerando embeddings e indexando {len(todos_documentos)} chunks no ChromaDB...")
    try:
        Chroma.from_documents(
            documents=todos_documentos,
            embedding=embeddings,
            persist_directory=str(PASTA_CHROMA),
        )
    except Exception as e:
        log_fn(f"Erro ao indexar documentos no ChromaDB: {e}")
        logger.ingestao(0, len(todos_documentos), False)
        return False

    log_fn(f"\nOK Concluído! {len(todos_documentos)} chunks de {arquivos_processados} arquivos indexados em {PASTA_CHROMA}/")
    log_fn(f"  Departamentos: {[p.name for p in pastas_departamento]}")
    logger.ingestao(arquivos_processados, len(todos_documentos), True)
    return True


if __name__ == "__main__":
    processar_documentos()
